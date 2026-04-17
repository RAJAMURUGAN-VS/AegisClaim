import os
import io
import re
from flask import Flask, request, jsonify
from flask_cors import CORS
from paddleocr import PaddleOCR
from PIL import Image
from pdf2image import convert_from_bytes
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation

app = Flask(__name__)
CORS(app)

# Initialize OCR
ocr = PaddleOCR(use_angle_cls=True, lang='en')

UPLOAD_FOLDER = 'uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# -------------------------------
# FILE TYPE DETECTION
# -------------------------------
def get_file_type(filename):
    ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
    if ext in ['png', 'jpg', 'jpeg', 'bmp', 'tiff', 'tif', 'webp', 'gif']:
        return 'image'
    elif ext == 'pdf':
        return 'pdf'
    elif ext == 'docx':
        return 'docx'
    elif ext in ['xlsx', 'xls']:
        return 'excel'
    elif ext == 'pptx':
        return 'pptx'
    else:
        return 'unknown'

# -------------------------------
# OCR CORE
# -------------------------------
def run_ocr_on_image(pil_image):
    import numpy as np

    img_array = np.array(pil_image.convert('RGB'))
    result = ocr.ocr(img_array)

    lines = []
    if result and result[0]:
        for line in result[0]:
            text = line[1][0]
            confidence = round(float(line[1][1]), 4)

            lines.append({
                'text': text,
                'confidence': confidence
            })

    return lines

# -------------------------------
# SMART FIELD EXTRACTION (IMPROVED)
# -------------------------------
def extract_fields(lines):

    # Join text
    text = " ".join([l['text'] for l in lines])

    # Normalize text (VERY IMPORTANT)
    text = re.sub(r'\s+', ' ', text)

    def find(pattern):
        match = re.search(pattern, text, re.IGNORECASE)
        return match.group(1).strip() if match else None

    return {
        "claim_number": find(r"Claim\s*Number\s*([A-Z0-9\-]+)"),

        "patient_name": find(
            r"Patient\s*Name\s*[:\-]?\s*([A-Z][a-z]+(?:\s[A-Z][a-z]+)+)"
        ),

        "policy_number": find(
            r"Policy\s*Number\s*[:\-]?\s*([A-Z0-9]+)"
        ),

        "date": find(
            r"Date\s*[:\-]?\s*(\d{2}/\d{2}/\d{4})"
        ),

        "amount": find(
            r"Total\s*Authorized\s*amount.*?(\d{5,})"
        ),

        "hospital": find(
            r"([A-Z]+\s*HOSPITAL)"
        ),

        "diagnosis": find(
            r"Diagnosis\s*[:\-]?\s*([A-Z]+)"
        )
    }

# -------------------------------
# PROCESSORS
# -------------------------------
def process_image(file_bytes, filename):
    try:
        image = Image.open(io.BytesIO(file_bytes))
        lines = run_ocr_on_image(image)

        return {
            'success': True,
            'file_type': 'image',
            'filename': filename,
            'lines': lines,
            'page_count': 1
        }

    except Exception as e:
        return {'success': False, 'error': str(e)}

def process_pdf(file_bytes, filename):
    try:
        pages = convert_from_bytes(file_bytes, dpi=200)

        all_lines = []

        for page_num, page_img in enumerate(pages, start=1):
            lines = run_ocr_on_image(page_img)

            for l in lines:
                l['page'] = page_num

            all_lines.extend(lines)

        return {
            'success': True,
            'file_type': 'pdf',
            'filename': filename,
            'lines': all_lines,
            'page_count': len(pages)
        }

    except Exception as e:
        return {'success': False, 'error': str(e)}

def process_docx(file_bytes, filename):
    try:
        doc = DocxDocument(io.BytesIO(file_bytes))

        lines = []
        for para in doc.paragraphs:
            if para.text.strip():
                lines.append({'text': para.text.strip(), 'confidence': 1.0})

        for table in doc.tables:
            for row in table.rows:
                for c in row.cells:
                    if c.text.strip():
                        lines.append({'text': c.text.strip(), 'confidence': 1.0})

        return {
            'success': True,
            'file_type': 'docx',
            'filename': filename,
            'lines': lines,
            'page_count': 1
        }

    except Exception as e:
        return {'success': False, 'error': str(e)}

def process_excel(file_bytes, filename):
    try:
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)

        lines = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                for val in row:
                    if val:
                        lines.append({'text': str(val), 'confidence': 1.0})

        return {
            'success': True,
            'file_type': 'excel',
            'filename': filename,
            'lines': lines,
            'page_count': len(wb.worksheets)
        }

    except Exception as e:
        return {'success': False, 'error': str(e)}

def process_pptx(file_bytes, filename):
    try:
        prs = Presentation(io.BytesIO(file_bytes))

        lines = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    lines.append({
                        'text': shape.text.strip(),
                        'confidence': 1.0,
                        'slide': slide_num
                    })

        return {
            'success': True,
            'file_type': 'pptx',
            'filename': filename,
            'lines': lines,
            'page_count': len(prs.slides)
        }

    except Exception as e:
        return {'success': False, 'error': str(e)}

# -------------------------------
# MAIN OCR API
# -------------------------------
@app.route('/ocr', methods=['POST'])
def extract_text():

    if 'file' not in request.files:
        return jsonify({'success': False, 'error': 'No file uploaded'}), 400

    uploaded_file = request.files['file']

    if uploaded_file.filename == '':
        return jsonify({'success': False, 'error': 'Empty filename'}), 400

    filename = uploaded_file.filename
    file_bytes = uploaded_file.read()
    file_type = get_file_type(filename)

    if file_type == 'image':
        result = process_image(file_bytes, filename)
    elif file_type == 'pdf':
        result = process_pdf(file_bytes, filename)
    elif file_type == 'docx':
        result = process_docx(file_bytes, filename)
    elif file_type == 'excel':
        result = process_excel(file_bytes, filename)
    elif file_type == 'pptx':
        result = process_pptx(file_bytes, filename)
    else:
        result = process_image(file_bytes, filename)

    # 🔥 CLEAN + EXTRACT
    if result.get("success"):

        # remove low-confidence garbage
        clean_lines = [
            l for l in result["lines"] if l["confidence"] > 0.6
        ]

        result["clean_lines"] = clean_lines
        result["extracted"] = extract_fields(clean_lines)

    return jsonify(result), 200 if result.get('success') else 500

# -------------------------------
# HEALTH CHECK
# -------------------------------
@app.route('/health', methods=['GET'])
def health_check():
    return jsonify({'status': 'running', 'ocr': 'ready'}), 200

# -------------------------------
# RUN SERVER
# -------------------------------
if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)