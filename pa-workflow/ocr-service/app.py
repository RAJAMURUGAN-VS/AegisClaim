import io
import os
import re
import tempfile

os.environ["PADDLE_PDX_DISABLE_MODEL_SOURCE_CHECK"] = "True"

from flask import Flask, jsonify, request
from flask_cors import CORS
from paddleocr import PaddleOCR
from pypdf import PdfReader
import fitz
from PIL import Image
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation

try:
    from pdf2image import convert_from_bytes
except ImportError:
    convert_from_bytes = None

app = Flask(__name__)
CORS(app)

ocr = PaddleOCR(use_textline_orientation=True, lang="en")


def get_file_type(filename: str) -> str:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext in ["png", "jpg", "jpeg", "bmp", "tiff", "tif", "webp", "gif"]:
        return "image"
    if ext == "pdf":
        return "pdf"
    if ext == "docx":
        return "docx"
    if ext in ["xlsx", "xls"]:
        return "excel"
    if ext == "pptx":
        return "pptx"
    return "unknown"


def run_ocr_on_image(pil_image: Image.Image):
    import numpy as np

    img_array = np.array(pil_image.convert("RGB"))
    result = ocr.ocr(img_array)

    lines = []
    if result and result[0]:
        for line in result[0]:
            text = line[1][0]
            confidence = round(float(line[1][1]), 4)
            lines.append({"text": text, "confidence": confidence})
    return lines


def extract_fields(lines):
    text = " ".join([line["text"] for line in lines])
    text = re.sub(r"\s+", " ", text)

    def find(pattern):
        match = re.search(pattern, text, re.IGNORECASE)
        return match.group(1).strip() if match else None

    return {
        "claim_number": find(r"Claim\s*Number\s*([A-Z0-9\-]+)"),
        "patient_name": find(r"Patient\s*Name\s*[:\-]?\s*([A-Z][a-z]+(?:\s[A-Z][a-z]+)+)"),
        "policy_number": find(r"Policy\s*Number\s*[:\-]?\s*([A-Z0-9]+)"),
        "date": find(r"Date\s*[:\-]?\s*(\d{2}/\d{2}/\d{4})"),
        "amount": find(r"Total\s*Authorized\s*amount.*?(\d{5,})"),
        "hospital": find(r"([A-Z]+\s*HOSPITAL)"),
        "diagnosis": find(r"Diagnosis\s*[:\-]?\s*([A-Z]+)"),
    }


def process_image(file_bytes, filename):
    try:
        image = Image.open(io.BytesIO(file_bytes))
        lines = run_ocr_on_image(image)
        return {
            "success": True,
            "file_type": "image",
            "filename": filename,
            "lines": lines,
            "page_count": 1,
        }
    except Exception as exc:
        return {"success": False, "error": str(exc)}


def process_pdf(file_bytes, filename):
    try:
        if convert_from_bytes is None:
            raise RuntimeError("pdf2image is not installed")

        pages = convert_from_bytes(file_bytes, dpi=200)
        all_lines = []
        for page_num, page_img in enumerate(pages, start=1):
            lines = run_ocr_on_image(page_img)
            for line in lines:
                line["page"] = page_num
            all_lines.extend(lines)

        return {
            "success": True,
            "file_type": "pdf",
            "filename": filename,
            "lines": all_lines,
            "page_count": len(pages),
        }
    except Exception as exc:
        # Fallback for environments where Poppler is not installed (common on Windows).
        error_message = str(exc)
        if "poppler" in error_message.lower() or "pdf2image" in error_message.lower():
            temp_pdf_path = None
            try:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
                    temp_pdf.write(file_bytes)
                    temp_pdf_path = temp_pdf.name

                # Pure-Python fallback: extract embedded text without Poppler.
                # This works for digital PDFs and avoids platform-level dependencies.
                reader = PdfReader(temp_pdf_path)
                all_lines = []

                for page_num, page in enumerate(reader.pages, start=1):
                    page_text = (page.extract_text() or "").strip()
                    if not page_text:
                        continue

                    for raw_line in page_text.splitlines():
                        clean_text = raw_line.strip()
                        if clean_text:
                            all_lines.append({"text": clean_text, "confidence": 1.0, "page": page_num})

                # If no embedded text is present, render pages with PyMuPDF and OCR the images.
                if not all_lines:
                    rendered_lines = []
                    pdf_doc = fitz.open(temp_pdf_path)
                    for page_num, page in enumerate(pdf_doc, start=1):
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                        page_lines = run_ocr_on_image(img)
                        for line in page_lines:
                            line["page"] = page_num
                        rendered_lines.extend(page_lines)
                    pdf_doc.close()
                    all_lines = rendered_lines

                return {
                    "success": True,
                    "file_type": "pdf",
                    "filename": filename,
                    "lines": all_lines,
                    "page_count": len(reader.pages),
                    "warning": "Poppler not found. Used pypdf/PyMuPDF fallback.",
                }
            except Exception as fallback_exc:
                return {"success": False, "error": f"{error_message}; fallback failed: {fallback_exc}"}
            finally:
                if temp_pdf_path and os.path.exists(temp_pdf_path):
                    os.remove(temp_pdf_path)

        return {"success": False, "error": error_message}


def process_docx(file_bytes, filename):
    try:
        doc = DocxDocument(io.BytesIO(file_bytes))
        lines = []

        for para in doc.paragraphs:
            if para.text.strip():
                lines.append({"text": para.text.strip(), "confidence": 1.0})

        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        lines.append({"text": cell.text.strip(), "confidence": 1.0})

        return {
            "success": True,
            "file_type": "docx",
            "filename": filename,
            "lines": lines,
            "page_count": 1,
        }
    except Exception as exc:
        return {"success": False, "error": str(exc)}


def process_excel(file_bytes, filename):
    try:
        workbook = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
        lines = []

        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                for value in row:
                    if value:
                        lines.append({"text": str(value), "confidence": 1.0})

        return {
            "success": True,
            "file_type": "excel",
            "filename": filename,
            "lines": lines,
            "page_count": len(workbook.worksheets),
        }
    except Exception as exc:
        return {"success": False, "error": str(exc)}


def process_pptx(file_bytes, filename):
    try:
        presentation = Presentation(io.BytesIO(file_bytes))
        lines = []

        for slide_num, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append({"text": shape.text.strip(), "confidence": 1.0, "slide": slide_num})

        return {
            "success": True,
            "file_type": "pptx",
            "filename": filename,
            "lines": lines,
            "page_count": len(presentation.slides),
        }
    except Exception as exc:
        return {"success": False, "error": str(exc)}


@app.post("/ocr")
def extract_text():
    if "file" not in request.files:
        return jsonify({"success": False, "error": "No file uploaded"}), 400

    uploaded_file = request.files["file"]
    if uploaded_file.filename == "":
        return jsonify({"success": False, "error": "Empty filename"}), 400

    filename = uploaded_file.filename
    file_bytes = uploaded_file.read()
    file_type = get_file_type(filename)

    if file_type == "image":
        result = process_image(file_bytes, filename)
    elif file_type == "pdf":
        result = process_pdf(file_bytes, filename)
    elif file_type == "docx":
        result = process_docx(file_bytes, filename)
    elif file_type == "excel":
        result = process_excel(file_bytes, filename)
    elif file_type == "pptx":
        result = process_pptx(file_bytes, filename)
    else:
        result = process_image(file_bytes, filename)

    if result.get("success"):
        clean_lines = [line for line in result["lines"] if line["confidence"] > 0.6]
        result["clean_lines"] = clean_lines
        result["extracted"] = extract_fields(clean_lines)
    else:
        app.logger.error("OCR failed for %s: %s", filename, result.get("error", "Unknown error"))

    return jsonify(result), 200 if result.get("success") else 500


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5001, debug=True, use_reloader=False)
